from flask import Flask, render_template, request, send_file
from openpyxl import load_workbook
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import copy
from pptx.dml.color import RGBColor
from PIL import Image

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['STATIC_FOLDER'] = 'static'

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/download_template')
def download_template():
    return send_file(os.path.join(app.config['STATIC_FOLDER'], 'template.xlsx'), as_attachment=True)

@app.route('/download_template_uk')
def download_template_uk():
    return send_file(os.path.join(app.config['STATIC_FOLDER'], 'template_uk.xlsx'), as_attachment=True)

@app.route('/download_explanation')
def download_explanation():
    return send_file(os.path.join(app.config['STATIC_FOLDER'], 'explanation.pptx'), as_attachment=True)

@app.route('/upload_files', methods=['POST'])
def upload_files():
    selected_script = request.form.get('script-select')
    if selected_script == 'Euro/Metric':
        # Handle uploaded data.xlsx file
        data_file = request.files['data_file']
        data_path = os.path.join(app.config['UPLOAD_FOLDER'], 'data.xlsx')
        data_file.save(data_path)
        
        # Handle uploaded picture files
        picture_paths = []
        for picture_file in request.files.getlist('picture_files'):
            picture_path = os.path.join(app.config['UPLOAD_FOLDER'], picture_file.filename)
            picture_file.save(picture_path)
            picture_paths.append(picture_path)
        
        # Load the Excel file
        workbook = load_workbook(data_path)
        worksheet = workbook.active
        
        # Loop through the cells in the worksheet
        empty_rows = []
        for i, row in enumerate(worksheet.iter_rows()):
            if all(cell.value is None for cell in row):
                empty_rows.append(i + 1)
        
        # Delete empty rows
        for row_index in reversed(empty_rows):
            worksheet.delete_rows(row_index)
            
        # Make None values empty string
        for i, row in enumerate(worksheet.iter_rows(min_row=3)):
            for cell in row:
                if cell.value is None:
                    cell.value = ""
        
        # Load the PowerPoint file
        ppt = Presentation(os.path.join(app.config['STATIC_FOLDER'], 'template.pptx'))
        
        # Get the first slide in the presentation
        slide = ppt.slides[0]
        
        # Loop through each row in the Excel file, starting with the second row
        for x in worksheet.iter_rows(min_row=3, values_only=True):
            # Duplicate the first slide
            copied_slide = ppt.slides.add_slide(slide.slide_layout)
        
            # Copy all the shapes from the original slide to the copied slide, skipping placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    continue
                el = shape.element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        
            # Title
            textbox = copied_slide.shapes[12]
            textbox.text = str(x[1])
            textbox.text_frame.paragraphs[0].font.name = 'Financier Display'
            textbox.text_frame.paragraphs[0].font.size = Pt(28)
        
            #Number
            textbox = copied_slide.shapes[11]
            textbox.text = str(x[0])
            textbox.text_frame.paragraphs[0].font.name = 'Financier Display'
            textbox.text_frame.paragraphs[0].font.size = Pt(28)
            
            #Table 4 - Shape 3 - Property Status
            table = copied_slide.shapes[3]
            cell = table.table.cell(0, 1)
            cell.text = str(x[2])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
        
            #Table 4 - Shape 3 - Date available
            table = copied_slide.shapes[3]
            cell = table.table.cell(1, 1)
            cell.text = str(x[3])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38) 
            
            #Table 4 - Shape 3 - construction start
            table = copied_slide.shapes[3]
            cell = table.table.cell(2, 1)
            cell.text = str(x[4])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38) 
            
            #Table 5 - Shape 4 - Plot
            table = copied_slide.shapes[4]
            cell = table.table.cell(0, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[5])))
            except:
                cell.text = str(x[5])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Warehouse
            table = copied_slide.shapes[4]
            cell = table.table.cell(1, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[6])))
            except:
                cell.text = str(x[6])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Office
            table = copied_slide.shapes[4]
            cell = table.table.cell(2, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[7])))
            except:
                cell.text = str(x[7])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Mezzanine
            table = copied_slide.shapes[4]
            cell = table.table.cell(3, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[8])))
            except:
                cell.text = str(x[8])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Clear Height
            table = copied_slide.shapes[5]
            cell = table.table.cell(0, 1)
            cell.text = '{} m'.format(str(x[9]))
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Floor load
            table = copied_slide.shapes[5]
            cell = table.table.cell(1, 1)
            cell.text = '{} tons/sq. m.'.format(str(x[10]))
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - sprinklers
            table = copied_slide.shapes[5]
            cell = table.table.cell(2, 1)
            cell.text = str(x[11])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Loading docks
            table = copied_slide.shapes[5]
            cell = table.table.cell(3, 1)
            cell.text = str(x[12])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - overhead doors
            table = copied_slide.shapes[5]
            cell = table.table.cell(4, 1)
            cell.text = str(x[13])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 - truck parcking
            table = copied_slide.shapes[9]
            cell = table.table.cell(0, 1)
            cell.text = str(x[14])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 - personal parking
            table = copied_slide.shapes[9]
            cell = table.table.cell(1, 1)
            cell.text = str(x[15])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 9 - Shape 8 - BREAAM
            table = copied_slide.shapes[8]
            cell = table.table.cell(0, 1)
            cell.text = str(x[16])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 1 - WH Size
            table = copied_slide.shapes[6]
            cell = table.table.cell(1, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[17])))
            except:
                cell.text = str(x[17])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 1 - Office
            table = copied_slide.shapes[6]
            cell = table.table.cell(1, 2)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[18])))
            except:
                cell.text = str(x[18])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 1 - Mezzanine
            table = copied_slide.shapes[6]
            cell = table.table.cell(1, 3)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[19])))
            except:
                cell.text = str(x[19])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 2 - WH Size
            table = copied_slide.shapes[6]
            cell = table.table.cell(2, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[20])))
            except:
                cell.text = str(x[20])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 2 - Office
            table = copied_slide.shapes[6]
            cell = table.table.cell(2, 2)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[21])))
            except:
                cell.text = str(x[21])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 2 - Mezzanine
            table = copied_slide.shapes[6]
            cell = table.table.cell(2, 3)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[22])))
            except:
                cell.text = str(x[22])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 3 - WH Size
            table = copied_slide.shapes[6]
            cell = table.table.cell(3, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[23])))
            except:
                cell.text = str(x[23])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 3 - Office
            table = copied_slide.shapes[6]
            cell = table.table.cell(3, 2)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[24])))
            except:
                cell.text = str(x[24])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 3 - Mezzanine
            table = copied_slide.shapes[6]
            cell = table.table.cell(3, 3)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[25])))
            except:
                cell.text = str(x[25])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 4 - WH Size
            table = copied_slide.shapes[6]
            cell = table.table.cell(4, 1)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[26])))
            except:
                cell.text = str(x[26])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 4 - Office
            table = copied_slide.shapes[6]
            cell = table.table.cell(4, 2)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[27])))
            except:
                cell.text = str(x[27])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 7 - Shape 6 - Unit 4 - Mezzanine
            table = copied_slide.shapes[6]
            cell = table.table.cell(4, 3)
            try:
                cell.text = '{:,.0f} sq. m.'.format(float(str(x[28])))
            except:
                cell.text = str(x[28])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 
            table = copied_slide.shapes[10]
            cell = table.table.cell(0, 1)
            try:
                cell.text = '??? {:,.2f} per sq. m. per annum'.format(float(str(x[29])))
            except:
                cell.text = str(x[29])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 
            table = copied_slide.shapes[10]
            cell = table.table.cell(1, 1)
            try:
                cell.text = '??? {:,.2f} per sq. m. per annum'.format(float(str(x[30])))
            except:
                cell.text = str(x[30])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 
            table = copied_slide.shapes[10]
            cell = table.table.cell(2, 1)
            try:
                cell.text = '??? {:,.2f} per sq. m. per annum'.format(float(str(x[31])))
            except:
                cell.text = str(x[31])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 
            table = copied_slide.shapes[10]
            cell = table.table.cell(3, 1)
            try:
                cell.text = '??? {:,.2f} per sq. m. per annum'.format(float(str(x[32])))
            except:
                cell.text = str(x[32])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 
            table = copied_slide.shapes[10]
            cell = table.table.cell(4, 1)
            try:
                cell.text = '??? {:,.2f} per sq. m. per annum'.format(float(str(x[33])))
            except:
                cell.text = str(x[33])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
        
            textframe = copied_slide.shapes[2]
            p = textframe.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = 'Google Maps'
            hlink = r.hyperlink
            hlink.address = x[36]
        
            #Picture placeholder 1
            table = copied_slide.shapes[0]
            image1_filename = str(x[34])
            extensions = ['.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG']
            for extension in extensions:
                image_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{image1_filename}{extension}")
                if os.path.isfile(image_path):
                    # Add the extension to the file path
                    break
            
            # Check if the file exists in the upload folder
            if os.path.isfile(image_path):
                # Insert the image into the table
                table = table.insert_picture(image_path)
            else:
                # Handle file not found in upload folder
                print(f"File not found in upload folder: {image1_filename}")  
            
            #Picture placeholder 2
            table = copied_slide.shapes[1]
            image2_filename = str(x[35])
            extensions = ['.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG']
            for extension in extensions:
                image2_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{image2_filename}{extension}")
                if os.path.isfile(image2_path):
                    # Add the extension to the file path
                    break
            
            # Check if the file exists in the upload folder
            if os.path.isfile(image2_path):
                # Insert the image into the table with the correct format
                table = table.insert_picture(image2_path)
            else:
                # Handle file not found in upload folder
                print(f"File not found in upload folder: {image2_filename}")
            
            
        # Save populated PowerPoint file
        ppt.save(os.path.join(app.config['UPLOAD_FOLDER'], 'mypopulated.pptx'))
        
        # Delete temporary data and picture files
        os.remove(data_path)
        for picture_path in picture_paths:
            os.remove(picture_path)
            
        # Return populated PowerPoint file for download
        return send_file(os.path.join(app.config['UPLOAD_FOLDER'], 'mypopulated.pptx'), as_attachment=True)
    elif selected_script == 'Pound/Imperial':
        data_file = request.files['data_file']
        data_path = os.path.join(app.config['UPLOAD_FOLDER'], 'data.xlsx')
        data_file.save(data_path)

        # Handle uploaded picture files
        picture_paths = []
        for picture_file in request.files.getlist('picture_files'):
            picture_path = os.path.join(app.config['UPLOAD_FOLDER'], picture_file.filename)
            picture_file.save(picture_path)
            picture_paths.append(picture_path)

        # Load the Excel file
        workbook = load_workbook(data_path)
        worksheet = workbook.active
        
        # Loop through the cells in the worksheet
        empty_rows = []
        for i, row in enumerate(worksheet.iter_rows()):
            if all(cell.value is None for cell in row):
                empty_rows.append(i + 1)

        # Delete empty rows
        for row_index in reversed(empty_rows):
            worksheet.delete_rows(row_index)
            
        # Make None values empty string
        for i, row in enumerate(worksheet.iter_rows(min_row=3)):
            for cell in row:
                if cell.value is None:
                    cell.value = ""

        # Load the PowerPoint file
        ppt = Presentation(os.path.join(app.config['STATIC_FOLDER'], 'template_uk.pptx'))
        
        # Get the first slide in the presentation
        slide = ppt.slides[0]
        
        # Loop through each row in the Excel file, starting with the second row
        for x in worksheet.iter_rows(min_row=3, values_only=True):
            # Duplicate the first slide
            copied_slide = ppt.slides.add_slide(slide.slide_layout)
        
            # Copy all the shapes from the original slide to the copied slide, skipping placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    continue
                el = shape.element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        
            # Title
            textbox = copied_slide.shapes[11]
            textbox.text = str(x[1])
            textbox.text_frame.paragraphs[0].font.name = 'Financier Display'
            textbox.text_frame.paragraphs[0].font.size = Pt(28)
        
            #Number
            textbox = copied_slide.shapes[10]
            textbox.text = str(x[0])
            textbox.text_frame.paragraphs[0].font.name = 'Financier Display'
            textbox.text_frame.paragraphs[0].font.size = Pt(28)
            
            #Table 4 - Shape 3 - Property Status
            table = copied_slide.shapes[3]
            cell = table.table.cell(0, 1)
            cell.text = str(x[2])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
        
            #Table 4 - Shape 3 - Date available
            table = copied_slide.shapes[3]
            cell = table.table.cell(1, 1)
            cell.text = str(x[3])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38) 
            
            #Table 4 - Shape 3 - construction start
            table = copied_slide.shapes[3]
            cell = table.table.cell(2, 1)
            cell.text = str(x[4])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38) 
            
            #Table 5 - Shape 4 - Plot
            table = copied_slide.shapes[4]
            cell = table.table.cell(0, 1)
            try:
                cell.text = '{:,.0f} sq. ft.'.format(float(str(x[5])))
            except:
                cell.text = str(x[5])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Warehouse
            table = copied_slide.shapes[4]
            cell = table.table.cell(1, 1)
            try:
                cell.text = '{:,.0f} sq. ft.'.format(float(str(x[6])))
            except:
                cell.text = str(x[6])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Office
            table = copied_slide.shapes[4]
            cell = table.table.cell(2, 1)
            try:
                cell.text = '{:,.0f} sq. ft.'.format(float(str(x[7])))
            except:
                cell.text = str(x[7])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 5 - Shape 4 - Mezzanine
            table = copied_slide.shapes[4]
            cell = table.table.cell(3, 1)
            try:
                cell.text = '{:,.0f} sq. ft.'.format(float(str(x[8])))
            except:
                cell.text = str(x[8])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Clear Height
            table = copied_slide.shapes[5]
            cell = table.table.cell(0, 1)
            cell.text = '{} ft'.format(str(x[9]))
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Floor load
            table = copied_slide.shapes[5]
            cell = table.table.cell(1, 1)
            cell.text = '{} lb/ft'.format(str(x[10]))
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - sprinklers
            table = copied_slide.shapes[5]
            cell = table.table.cell(2, 1)
            cell.text = str(x[11])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - Loading docks
            table = copied_slide.shapes[5]
            cell = table.table.cell(3, 1)
            cell.text = str(x[12])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 6 - Shape 5 - overhead doors
            table = copied_slide.shapes[5]
            cell = table.table.cell(4, 1)
            cell.text = str(x[13])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 - truck parcking
            table = copied_slide.shapes[8]
            cell = table.table.cell(0, 1)
            cell.text = str(x[14])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 10 - Shape 9 - personal parking
            table = copied_slide.shapes[8]
            cell = table.table.cell(1, 1)
            cell.text = str(x[15])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 9 - Shape 8 - BREAAM
            table = copied_slide.shapes[7]
            cell = table.table.cell(0, 1)
            cell.text = str(x[16])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 12 - Shape 9 - rental price
            table = copied_slide.shapes[9]
            cell = table.table.cell(0, 1)
            try:
                cell.text = '?? {:,.2f} per sq. ft. per annum'.format(float(str(x[17])))
            except:
                cell.text = str(x[17])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 9 - Shape 6 - Description
            table = copied_slide.shapes[6]
            cell = table.table.cell(0, 0)
            cell.text = str(x[20])
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 3 - Shape 12 - Owner developer
            table = copied_slide.shapes[12]
            cell = table.table.cell(0, 1)
            cell.text = str(x[18])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
            
            #Table 3 - Shape 12 - Planning
            table = copied_slide.shapes[12]
            cell = table.table.cell(1, 1)
            cell.text = str(x[19])
            cell.text_frame.paragraphs[0].font.name = 'Calibre'
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT   
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(38, 38, 38)
       
            textframe = copied_slide.shapes[2]
            p = textframe.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = 'Google Maps'
            hlink = r.hyperlink
            hlink.address = x[23]
        
            #Picture placeholder 1
            table = copied_slide.shapes[0]
            image1_filename = str(x[21])
            extensions = ['.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG']
            for extension in extensions:
                image_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{image1_filename}{extension}")
                if os.path.isfile(image_path):
                    # Add the extension to the file path
                    break
            
            # Check if the file exists in the upload folder
            if os.path.isfile(image_path):
                # Insert the image into the table
                table = table.insert_picture(image_path)
            else:
                # Handle file not found in upload folder
                print(f"File not found in upload folder: {image1_filename}")  
            
            #Picture placeholder 2
            table = copied_slide.shapes[1]
            image2_filename = str(x[22])
            extensions = ['.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG']
            for extension in extensions:
                image2_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{image2_filename}{extension}")
                if os.path.isfile(image2_path):
                    # Add the extension to the file path
                    break
            
            # Check if the file exists in the upload folder
            if os.path.isfile(image2_path):
                # Insert the image into the table with the correct format
                table = table.insert_picture(image2_path)
            else:
                # Handle file not found in upload folder
                print(f"File not found in upload folder: {image2_filename}")
            
            
        # Save populated PowerPoint file
        ppt.save(os.path.join(app.config['UPLOAD_FOLDER'], 'mypopulated.pptx'))
        
        # Delete temporary data and picture files
        os.remove(data_path)
        for picture_path in picture_paths:
            os.remove(picture_path)
            
        # Return populated PowerPoint file for download
        return send_file(os.path.join(app.config['UPLOAD_FOLDER'], 'mypopulated.pptx'), as_attachment=True)
    else:
        return 'Invalid script selection'
    
if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
